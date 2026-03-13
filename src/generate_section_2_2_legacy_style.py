from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import sys
import textwrap

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from matplotlib import font_manager
from matplotlib.ticker import PercentFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

BASE_DIR = Path(__file__).resolve().parent.parent
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from src.generate_section_2_2_assets import SCENE_COLORS, SCENE_ORDER, load_bundle


DOCX_PATH = BASE_DIR / "Section_2.2_应用现状综述_最终版.docx"
PPTX_PATH = BASE_DIR / "Section_2.2_应用现状综述_汇报版.pptx"
NOTES_PATH = BASE_DIR / "Section_2.2_应用现状综述_汇报讲稿.md"
PPT_FIG_DIR = BASE_DIR / "output" / "figures" / "section_2_2_ppt_v6"

LEGACY_FIGURES = {
    "province": BASE_DIR / "output" / "figures" / "fig_a01_province_map.png",
    "stage": BASE_DIR / "output" / "figures" / "fig_a02_stage_waffle.png",
    "subject": BASE_DIR / "output" / "figures" / "fig_a03_subject_lollipop.png",
    "scenario": BASE_DIR / "output" / "figures" / "fig_a04_scenario_treemap.png",
    "region_product": BASE_DIR / "output" / "figures" / "fig_a06_region_product.png",
    "stage_product": BASE_DIR / "output" / "figures" / "fig_a07_stage_product.png",
    "subject_product": BASE_DIR / "output" / "figures" / "fig_a08_subject_product.png",
}

PPT_THEME = {
    "bg": RGBColor(246, 241, 232),
    "ink": RGBColor(20, 42, 56),
    "muted": RGBColor(86, 103, 117),
    "panel": RGBColor(255, 252, 247),
    "accent": RGBColor(199, 92, 56),
    "olive": RGBColor(114, 126, 84),
    "line": RGBColor(225, 216, 203),
}

PROVINCE_REGION = {
    "北京市": "东部",
    "天津市": "东部",
    "河北省": "东部",
    "上海市": "东部",
    "江苏省": "东部",
    "浙江省": "东部",
    "福建省": "东部",
    "山东省": "东部",
    "广东省": "东部",
    "海南省": "东部",
    "辽宁省": "东部",
    "山西省": "中部",
    "吉林省": "中部",
    "黑龙江省": "中部",
    "安徽省": "中部",
    "江西省": "中部",
    "河南省": "中部",
    "湖北省": "中部",
    "湖南省": "中部",
    "内蒙古自治区": "西部",
    "广西壮族自治区": "西部",
    "重庆市": "西部",
    "四川省": "西部",
    "贵州省": "西部",
    "云南省": "西部",
    "西藏自治区": "西部",
    "陕西省": "西部",
    "甘肃省": "西部",
    "青海省": "西部",
    "宁夏回族自治区": "西部",
    "新疆维吾尔自治区": "西部",
}


@dataclass
class LegacyBundle:
    bundle: object
    df: pd.DataFrame
    cases: pd.DataFrame
    product_col: str
    region_product: pd.DataFrame
    stage_product: pd.DataFrame
    subject_product: pd.DataFrame
    scene_product: pd.DataFrame
    stage_scene_pct: pd.DataFrame
    region_scene_pct: pd.DataFrame
    scene_profiles: pd.DataFrame


def configure_matplotlib() -> None:
    preferred_font_names = [
        "Arial Unicode MS",
        "Hiragino Sans GB",
        "Songti SC",
        "Heiti TC",
        "PingFang HK",
        "DejaVu Sans",
    ]
    available_names = {font.name for font in font_manager.fontManager.ttflist}
    selected_font = next((name for name in preferred_font_names if name in available_names), "DejaVu Sans")
    plt.rcParams["font.family"] = [selected_font]
    plt.rcParams["font.sans-serif"] = [selected_font]
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["figure.facecolor"] = "white"
    sns.set_theme(style="whitegrid", rc={"font.family": selected_font, "axes.unicode_minus": False})


def normalize_stage_label(value: str) -> str:
    text = str(value or "").strip()
    if not text:
        return "未提及"
    if "幼儿园" in text or "学前" in text:
        return "幼儿园"
    if "小学" in text:
        return "小学"
    if "初中" in text:
        return "初中"
    if "高中" in text:
        return "高中"
    return "其他"


def pct(part: float, whole: float) -> str:
    if not whole:
        return "0.0%"
    return f"{part / whole * 100:.1f}%"


def list_phrase(items: list[str], limit: int = 4) -> str:
    items = [item for item in items if item][:limit]
    if not items:
        return "暂无明显头部产品"
    if len(items) == 1:
        return items[0]
    return "、".join(items[:-1]) + "和" + items[-1]


def top_items_text(series: pd.Series, limit: int = 3) -> str:
    items = [(str(name), int(value)) for name, value in series.head(limit).items() if int(value) > 0]
    return "、".join(f"{name}（{value}次）" for name, value in items)


def wrap_label(text: str, width: int = 10) -> str:
    return "\n".join(textwrap.wrap(str(text), width=width))


def ensure_legacy_figures() -> None:
    missing = [path for path in LEGACY_FIGURES.values() if not path.exists()]
    if missing:
        names = "\n".join(str(path) for path in missing)
        raise FileNotFoundError(f"缺少旧版图表文件，无法按旧风格生成 Word：\n{names}")


def regenerate_legacy_product_heatmaps(data: LegacyBundle) -> None:
    heatmap_specs = [
        (
            data.region_product,
            LEGACY_FIGURES["region_product"],
            "图5 典型区域与头部AI产品联动应用热力图",
            "AI产品",
            "区域",
            "Blues",
            (10, 4.2),
        ),
        (
            data.stage_product,
            LEGACY_FIGURES["stage_product"],
            "图6 学段与头部AI产品联动应用热力图",
            "AI产品",
            "学段",
            "Oranges",
            (10, 4.2),
        ),
        (
            data.subject_product,
            LEGACY_FIGURES["subject_product"],
            "图7 核心学科与头部AI产品联动应用热力图",
            "AI产品",
            "学科",
            "Greens",
            (10, 5),
        ),
    ]
    for frame, output, title, xlabel, ylabel, cmap, figsize in heatmap_specs:
        fig, ax = plt.subplots(figsize=figsize)
        sns.heatmap(frame, annot=True, fmt="g", cmap=cmap, linewidths=0.5, cbar=False, ax=ax)
        ax.set_title(title, pad=15, fontsize=14, fontweight="bold")
        ax.set_xlabel(xlabel, fontsize=12)
        ax.set_ylabel(ylabel, fontsize=12)
        ax.set_xticklabels([wrap_label(label.get_text(), 7) for label in ax.get_xticklabels()], rotation=0, ha="center", fontsize=9)
        ax.set_yticklabels([label.get_text() for label in ax.get_yticklabels()], rotation=0, fontsize=10)
        fig.tight_layout()
        fig.savefig(output, dpi=300)
        plt.close(fig)


def build_legacy_bundle() -> LegacyBundle:
    source = load_bundle()
    df = source.df.copy()
    cases = source.cases.copy()
    product_col = "产品名_标准"
    df["学段_旧版"] = df["学段"].map(normalize_stage_label)
    df["区域"] = df["省份_标准"].map(PROVINCE_REGION).fillna("未提及")

    top_products = df[product_col].dropna().value_counts().head(8).index.tolist()
    top_products_extended = df[product_col].dropna().value_counts().head(10).index.tolist()

    region_product = pd.crosstab(df["区域"], df[product_col]).reindex(index=["东部", "中部", "西部"], fill_value=0)
    region_product = region_product.reindex(columns=top_products, fill_value=0)

    stage_product = pd.crosstab(df["学段_旧版"], df[product_col]).reindex(index=["幼儿园", "小学", "初中", "高中"], fill_value=0)
    stage_product = stage_product.reindex(columns=top_products, fill_value=0)

    subject_order = [subject for subject in ["语文", "数学", "英语", "科学", "美术"] if subject in df["学科"].dropna().unique()]
    subject_product = pd.crosstab(df["学科"], df[product_col]).reindex(index=subject_order, fill_value=0)
    subject_product = subject_product.reindex(columns=top_products, fill_value=0)

    scene_product = pd.crosstab(cases["应用场景（一级）"], cases[product_col]).reindex(index=SCENE_ORDER, fill_value=0)
    scene_product = scene_product.reindex(columns=top_products_extended, fill_value=0)

    stage_scene_pct = pd.crosstab(cases["学段_标准"], cases["应用场景（一级）"], normalize="columns").mul(100)
    stage_scene_pct = stage_scene_pct.reindex(index=["学前", "小学", "初中", "高中"], columns=SCENE_ORDER, fill_value=0)

    region_scene_pct = pd.crosstab(cases["区域"], cases["应用场景（一级）"], normalize="index").mul(100)
    region_scene_pct = region_scene_pct.reindex(index=["东部", "中部", "西部"], columns=SCENE_ORDER, fill_value=0)

    scene_profiles = (
        cases.dropna(subset=["应用场景（一级）"])
        .groupby("应用场景（一级）")
        .agg(
            case_count=("案例编号", "nunique"),
            province_coverage=("省份_标准", lambda series: series.dropna().nunique()),
            subject_coverage=("学科", lambda series: series.dropna().nunique()),
            stage_coverage=("学段_标准", lambda series: series.dropna().nunique()),
        )
        .reindex(SCENE_ORDER)
        .fillna(0)
    )
    scene_profiles["top_stage"] = [
        pd.crosstab(cases["学段_标准"], cases["应用场景（一级）"])[scene].idxmax() if scene in cases["应用场景（一级）"].values else "未提及"
        for scene in SCENE_ORDER
    ]
    scene_profiles["top_subject"] = "未提及"
    subject_scene = pd.crosstab(cases["学科"].fillna("未提及"), cases["应用场景（一级）"])
    for scene in SCENE_ORDER:
        if scene in subject_scene.columns and subject_scene[scene].sum() > 0:
            scene_profiles.loc[scene, "top_subject"] = subject_scene[scene].sort_values(ascending=False).index[0]
        else:
            scene_profiles.loc[scene, "top_subject"] = "未提及"

    return LegacyBundle(
        bundle=source,
        df=df,
        cases=cases,
        product_col=product_col,
        region_product=region_product,
        stage_product=stage_product,
        subject_product=subject_product,
        scene_product=scene_product,
        stage_scene_pct=stage_scene_pct,
        region_scene_pct=region_scene_pct,
        scene_profiles=scene_profiles,
    )


def build_ppt_figures(data: LegacyBundle) -> dict[str, Path]:
    PPT_FIG_DIR.mkdir(parents=True, exist_ok=True)
    paths = {
        "scene_bubble": PPT_FIG_DIR / "fig_ppt_scene_bubble.png",
        "region_structure": PPT_FIG_DIR / "fig_ppt_region_structure.png",
        "stage_scene": PPT_FIG_DIR / "fig_ppt_stage_scene_pct.png",
        "scene_product": PPT_FIG_DIR / "fig_ppt_scene_product.png",
    }

    scene_profiles = data.scene_profiles.copy()
    fig, ax = plt.subplots(figsize=(9.6, 6.2))
    for scene in SCENE_ORDER:
        row = scene_profiles.loc[scene]
        ax.scatter(
            row["province_coverage"],
            row["subject_coverage"],
            s=max(row["case_count"], 1) * 1.4,
            color=SCENE_COLORS[scene],
            alpha=0.82,
            edgecolors="white",
            linewidths=1.5,
        )
        ax.text(
            row["province_coverage"],
            row["subject_coverage"] + 0.45,
            f"{scene}\n{int(row['case_count'])}例",
            ha="center",
            va="bottom",
            fontsize=10,
            weight="bold",
        )
    ax.set_title("六大场景的扩散广度与复杂度矩阵", fontsize=18, weight="bold")
    ax.set_xlabel("覆盖省份数")
    ax.set_ylabel("涉及学科数")
    ax.axvline(scene_profiles["province_coverage"].median(), linestyle="--", color="#A0A0A0", linewidth=1)
    ax.axhline(scene_profiles["subject_coverage"].median(), linestyle="--", color="#A0A0A0", linewidth=1)
    ax.text(0.03, 0.95, "右上角=既有规模也有跨学科扩散", transform=ax.transAxes, fontsize=10, color="#555555")
    sns.despine()
    fig.tight_layout()
    fig.savefig(paths["scene_bubble"], dpi=240)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(9.6, 5.3))
    bottom = pd.Series(0, index=data.region_scene_pct.index, dtype=float)
    for scene in SCENE_ORDER:
        values = data.region_scene_pct[scene]
        bars = ax.bar(data.region_scene_pct.index, values, bottom=bottom, color=SCENE_COLORS[scene], label=scene, width=0.62)
        for bar, value, baseline in zip(bars, values, bottom):
            if value >= 12:
                ax.text(bar.get_x() + bar.get_width() / 2, baseline + value / 2, f"{value:.0f}%", ha="center", va="center", fontsize=9, color="white", weight="bold")
        bottom += values
    ax.set_ylim(0, 100)
    ax.yaxis.set_major_formatter(PercentFormatter(100))
    ax.set_title("区域内部的场景结构差异", fontsize=18, weight="bold")
    ax.set_ylabel("区域内案例占比")
    ax.legend(ncol=6, bbox_to_anchor=(0.5, 1.16), loc="upper center", frameon=False)
    sns.despine()
    fig.tight_layout()
    fig.savefig(paths["region_structure"], dpi=240)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(8.8, 5.3))
    sns.heatmap(data.stage_scene_pct, annot=True, fmt=".0f", cmap="YlOrBr", linewidths=0.6, cbar_kws={"label": "场景内部占比(%)"}, ax=ax)
    ax.set_title("场景对学段的依赖程度", fontsize=18, weight="bold")
    ax.set_xlabel("")
    ax.set_ylabel("")
    fig.tight_layout()
    fig.savefig(paths["stage_scene"], dpi=240)
    plt.close(fig)

    scene_product_pct = data.scene_product.div(data.scene_product.sum(axis=1).replace(0, 1), axis=0).mul(100)
    fig, ax = plt.subplots(figsize=(11, 4.8))
    sns.heatmap(scene_product_pct, annot=True, fmt=".0f", cmap="crest", linewidths=0.6, cbar_kws={"label": "场景内产品占比(%)"}, ax=ax)
    ax.set_title("六大场景的头部产品控制力", fontsize=18, weight="bold")
    ax.set_xlabel("")
    ax.set_ylabel("")
    ax.set_xticklabels([wrap_label(label.get_text(), 6) for label in ax.get_xticklabels()], fontsize=9)
    fig.tight_layout()
    fig.savefig(paths["scene_product"], dpi=240)
    plt.close(fig)

    return paths


def set_docx_style(doc: Document) -> None:
    style = doc.styles["Normal"]
    style.font.name = "宋体"
    style._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    style.font.size = Pt(11)


def add_caption(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(text)
    run.font.size = Pt(10)


def insert_picture(doc: Document, path: Path, width: float = 5.0) -> None:
    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.add_run().add_picture(str(path), width=Inches(width))


def add_bold_lead(doc: Document, lead: str, tail: str) -> None:
    paragraph = doc.add_paragraph()
    lead_run = paragraph.add_run(lead)
    lead_run.bold = True
    paragraph.add_run(tail)


def build_word_doc(data: LegacyBundle) -> None:
    doc = Document()
    set_docx_style(doc)

    bundle = data.bundle
    cases = data.cases
    df = data.df
    summary = bundle.summary
    total_cases = summary["case_count"]
    total_rows = summary["tool_rows"]
    top_products = df[data.product_col].dropna().value_counts()
    top_stage = cases["学段_标准"].value_counts().index[0]
    top_stage_count = int(cases["学段_标准"].value_counts().iloc[0])
    top_subject_counts = cases["学科"].fillna("未提及").value_counts()
    main_subject = next((item for item in top_subject_counts.index if item != "未提及"), "语文")
    main_subject_count = int(top_subject_counts[main_subject])
    top_provinces = cases["省份_标准"].dropna().value_counts().head(4)
    top4_province_share = top_provinces.sum() / total_cases if total_cases else 0
    top_scene = pd.Series(summary["scene_counts"]).reindex(SCENE_ORDER).idxmax()
    top_scene_count = int(pd.Series(summary["scene_counts"]).reindex(SCENE_ORDER).max())
    scene_ranked = pd.Series(summary["scene_counts"]).sort_values(ascending=False)
    top2_scene_share = scene_ranked.head(2).sum() / scene_ranked.sum() if scene_ranked.sum() else 0
    east_cases = int((cases["区域"] == "东部").sum())
    middle_cases = int((cases["区域"] == "中部").sum())
    west_cases = int((cases["区域"] == "西部").sum())

    east_top = data.region_product.loc["东部"].sort_values(ascending=False)
    middle_top = data.region_product.loc["中部"].sort_values(ascending=False)
    west_top = data.region_product.loc["西部"].sort_values(ascending=False)
    stage_top = {stage: data.stage_product.loc[stage].sort_values(ascending=False) for stage in data.stage_product.index}
    subject_top = {subject: data.subject_product.loc[subject].sort_values(ascending=False) for subject in data.subject_product.index}
    scene_top = {scene: data.scene_product.loc[scene].sort_values(ascending=False) for scene in data.scene_product.index}

    doc.add_heading("2.2 教育应用类产品应用现状综述", level=1)
    doc.add_paragraph(
        f"当前，人工智能赋能基础教育正经历从单点探索向系统性应用的跨越。基于 V6.xlsx 主事实表“（新）processed_results”的更新数据，本节仍沿用原有 Section 2.2 的分析框架与图式，对 {total_cases} 个典型案例、{total_rows} 条工具/产品记录进行重新统计，并从学段学科渗透、区域空间分布与应用场景结构三个核心维度，全景式呈现 AI 教育应用的发展现状。"
    )

    doc.add_heading("2.2.1 学段学科角度：应用现状感知", level=2)
    doc.add_paragraph("AI教育应用在不同学段与学科间的渗透呈现出显著的结构性差异与聚集特征（见图2、图3）：")
    add_bold_lead(
        doc,
        "第一，通用基座模型的全域覆盖能力凸显。",
        f" 以{top_products.index[0]}为代表的头部通用模型，连同{list_phrase(top_products.index.tolist()[1:5], 4)}等多模态或平台型产品，在不同学段与学科中都呈现出较高覆盖率。这表明，具备自然语言理解、内容生成与教学资源重组能力的通用模型，已经成为当前教育创新实践的底层基础设施。",
    )
    add_bold_lead(
        doc,
        "第二，小学阶段成为应用落地的核心“试验田”。",
        f" 案例级统计显示，{top_stage}阶段共有{top_stage_count}例，占全部案例的{pct(top_stage_count, total_cases)}。这种占比优势与华夫饼图中的视觉主导地位相一致，也意味着在评价环境相对灵活、课堂创新空间较大的阶段，AI 工具更容易被纳入常态教学。",
    )
    add_bold_lead(
        doc,
        "第三，不同产品形态的学段适配性呈现分化。",
        doc_text_stage_tail(stage_top),
    )
    insert_picture(doc, LEGACY_FIGURES["stage"])
    add_caption(doc, "图2 学段分布华夫饼图")
    insert_picture(doc, LEGACY_FIGURES["subject"])
    add_caption(doc, "图3 学科渗透棒棒糖图")

    doc.add_heading("2.2.2 区域角度：空间分布格局", level=2)
    doc.add_paragraph("从地理空间视角审视（见图1），AI教育应用的区域分布与地方数字经济发展水平及教育信息化基础高度同频，呈现出典型的梯队差异：")
    add_bold_lead(
        doc,
        "第一，东部极化与“生态构建阶段”到来。",
        f" 在案例省域分布地图中，深色高频区域仍高度集中，{list_phrase(top_provinces.index.tolist(), 4)}构成当前最明显的头部省域板块，四省市合计占全部案例的{pct(top4_province_share, 1)}。东部地区样本达到{east_cases}例，显著高于中部的{middle_cases}例和西部的{west_cases}例，说明领先地区已经从单点试用转向平台、模型与教学流程的复合整合。",
    )
    add_bold_lead(
        doc,
        "第二，“多产品叠加”与“集群效应”显现。",
        f" 从区域×产品联动结果看，东部地区最活跃的头部产品为{top_items_text(east_top, 3)}；中部地区则由{top_items_text(middle_top, 3)}支撑其常态化教学数字化需求。这意味着领先区域不只是总量更多，更重要的是产品结构更丰富、组合关系更稳定。",
    )
    add_bold_lead(
        doc,
        "第三，西部地区的单点突破与规模困境。",
        f" 相比之下，西部地区的样本更多表现为头部产品牵引下的集中渗透，当前主要由{top_items_text(west_top, 3)}构成骨架。其特点不是完全缺位，而是更依赖少数低门槛、通用性强的产品完成快速接入，规模扩散能力仍有待提升。",
    )
    insert_picture(doc, LEGACY_FIGURES["province"])
    add_caption(doc, "图1 案例省域分布地图")
    doc.add_paragraph(
        "综合来看，当前区域差异的核心已经不只是案例总量高低，而是区域内部是否形成了稳定的产品组合与协同应用路径。东部更接近生态化整合阶段，中部处于平台化稳步推进阶段，西部则仍以头部通用产品牵引的导入式应用为主。"
    )

    doc.add_heading("2.2.3 场景角度：核心应用结构", level=2)
    doc.add_paragraph("对案例应用场景的分类剖析揭示了当前AI赋能教育的深层次偏向（见图4）：")
    add_bold_lead(
        doc,
        "首先，教育智能应用呈现显著的“助学主导”结构。",
        f" V6 案例中，{top_scene}场景共有{top_scene_count}例，占全部案例的{pct(top_scene_count, total_cases)}；前两大场景合计占比达到{pct(top2_scene_share, 1)}。在应用场景树图中，它仍然以压倒性的面积占据画面核心，说明当前 AI 教育应用并不是平均铺开，而是已经形成非常明确的主航道与次主航道。",
    )
    add_bold_lead(
        doc,
        "其次，“情境构建”与“智能辅导”成为核心路径。",
        f" 从二级场景看，{list_phrase(bundle.summary['top_l2_by_l1'].get('助学', []), 2)}依然构成助学场景的主轴；其背后对应的产品结构，则以{top_items_text(scene_top.get('助学', pd.Series(dtype=float)), 3)}为代表。换言之，头部模型正在通过“资源生成 + 个性化辅导”双路径嵌入课堂。",
    )
    add_bold_lead(
        doc,
        "最后，教研场景的智能化渗透依然薄弱。",
        f" 尽管助教、助评、助育、助管已形成更清晰的功能分层，但助研类案例仍然位于长尾位置，仅表现为{top_items_text(scene_top.get('助研', pd.Series(dtype=float)), 2)}等有限产品的探索性介入。教育评价、教学研究与制度性变革等复杂环节，依然需要更强的教育垂类模型与组织协同。",
    )
    insert_picture(doc, LEGACY_FIGURES["scenario"])
    add_caption(doc, "图4 应用场景树图")

    doc.add_heading("2.2.4 三维交叉：典型产品联动分析", level=2)
    doc.add_paragraph("为了进一步揭示头部AI产品在不同维度的真实渗透率与适配特征，本节基于应用产品频次映射，从三大维度进行深度联动交叉验证。")

    doc.add_paragraph("1. 区域 X 产品联动分析", style="List Number")
    add_region_bullet(doc, "东部地区", east_top, "呈现出头部模型与多模态工具并进的复合结构。")
    add_region_bullet(doc, "中部地区", middle_top, "以通用大模型为主导，但常态化教学平台的重要性更加突出。")
    add_region_bullet(doc, "西部地区", west_top, "应用生态呈现向头部产品集聚的“长尾收缩”特征。")
    insert_picture(doc, LEGACY_FIGURES["region_product"])
    add_caption(doc, "图5 典型区域与头部AI产品联动应用热力图")

    doc.add_paragraph("2. 学段 X 产品联动分析", style="List Number")
    for stage_name in ["幼儿园", "小学", "初中", "高中"]:
        if stage_name in stage_top:
            add_stage_bullet(doc, stage_name, stage_top[stage_name])
    insert_picture(doc, LEGACY_FIGURES["stage_product"])
    add_caption(doc, "图6 学段与头部AI产品联动应用热力图")

    doc.add_paragraph("3. 学科 X 产品联动分析", style="List Number")
    for subject_name in ["语文", "数学", "英语", "科学", "美术"]:
        if subject_name in subject_top:
            add_subject_bullet(doc, subject_name, subject_top[subject_name])
    insert_picture(doc, LEGACY_FIGURES["subject_product"])
    add_caption(doc, "图7 核心学科与头部AI产品联动应用热力图")
    doc.add_paragraph(
        "三维交叉结果表明，真正值得关注的已不是某个产品是否高频出现，而是它是否在特定区域、特定学段、特定学科中形成了稳定的联动优势。换言之，下一阶段竞争的核心不是通用热度，而是教育场景控制力。"
    )

    doc.add_heading("数据与代码使用说明（供核查参考）", level=3)
    doc.add_paragraph("为确保各项统计分析客观、精准溯源，本章节以上统计逻辑与作图说明如下：")
    add_bold_lead(doc, "使用数据源：", f" 本次正文已由旧版 V5 口径切换为 data/V6.xlsx，其中主事实表“（新）processed_results”包含{total_rows}条工具级记录，可回收为{total_cases}个案例。正文沿用旧版图组 A01-A08 的表达方式，但统计口径已更新为 V6。")
    add_bold_lead(doc, "本文交叉数据提取：", " 学段、学科、区域和产品联动结果均基于 V6 主事实表重新汇总；区域口径采用省份归一化映射，案例分布与产品频次采用案例级与工具级两套口径分别处理。")
    add_bold_lead(doc, "本文引用图表绘图代码：", " 旧版 A01-A08 图组继续使用项目中既有图件；本次 Section 2.2 的 Word 重建与洞见型 PPT/讲稿输出入口为 src/generate_section_2_2_legacy_style.py。")

    doc.save(DOCX_PATH)


def doc_text_stage_tail(stage_top: dict[str, pd.Series]) -> str:
    primary = stage_top.get("小学", pd.Series(dtype=float))
    high = stage_top.get("高中", pd.Series(dtype=float))
    kindergarten = stage_top.get("幼儿园", pd.Series(dtype=float))
    return (
        f" 工具型与内容创作型产品在低学段中的活跃度仍然更高，例如小学阶段主要由{top_items_text(primary, 3)}构成核心组合；"
        f"高中阶段则更明显地向{top_items_text(high, 2)}等面向任务完成与课程生产的产品收束；"
        f"而幼儿园/学前场景中，{top_items_text(kindergarten, 2)}等低门槛平台更容易进入教学实践，体现出学段差异背后的产品适配逻辑。"
    )


def add_region_bullet(doc: Document, region_name: str, series: pd.Series, suffix: str) -> None:
    paragraph = doc.add_paragraph(style="List Bullet")
    run = paragraph.add_run(f"{region_name}：")
    run.bold = True
    paragraph.add_run(f"{suffix} 当前最主要的产品组合为{top_items_text(series, 3)}。")


def add_stage_bullet(doc: Document, stage_name: str, series: pd.Series) -> None:
    stage_prefix = {
        "幼儿园": "应用生态显著区别于中小学。",
        "小学": "仍然是核心创新试验田。",
        "初中": "延续了小学的头部模型驱动，但更加聚焦刚性教学任务。",
        "高中": "应用结构出现向课程生产与任务完成倾斜的迹象。",
    }
    paragraph = doc.add_paragraph(style="List Bullet")
    run = paragraph.add_run(f"{stage_name}：")
    run.bold = True
    paragraph.add_run(f" {stage_prefix.get(stage_name, '')} 当前高频产品为{top_items_text(series, 3)}。")


def add_subject_bullet(doc: Document, subject_name: str, series: pd.Series) -> None:
    subject_prefix = {
        "语文": "语言理解与内容生成的深度融合是主旋律。",
        "数学": "逻辑推理与解题支持类模型优势更突出。",
        "英语": "语言训练、情境构建与作业反馈相互叠加。",
        "科学": "兼具知识解释与探究性任务支持特征。",
        "美术": "图像生成和创意构思工具占据明显优势。",
    }
    paragraph = doc.add_paragraph(style="List Bullet")
    run = paragraph.add_run(f"{subject_name}：")
    run.bold = True
    paragraph.add_run(f" {subject_prefix.get(subject_name, '')} 当前高频产品为{top_items_text(series, 3)}。")


def add_background(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptInches(13.333), PptInches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_panel(slide, left: float, top: float, width: float, height: float, fill: RGBColor | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or PPT_THEME["panel"]
    shape.line.color.rgb = PPT_THEME["line"]
    shape.line.width = PptPt(1)


def add_text(slide, left: float, top: float, width: float, height: float, text: str, size: int = 18, bold: bool = False, color: RGBColor | None = None, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "PingFang SC"
    paragraph.font.size = PptPt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color or PPT_THEME["ink"]


def add_metric_card(slide, left: float, top: float, title: str, value: str, accent: RGBColor) -> None:
    add_panel(slide, left, top, 2.0, 1.1)
    add_text(slide, left + 0.16, top + 0.16, 1.68, 0.28, title, 11, False, PPT_THEME["muted"])
    add_text(slide, left + 0.16, top + 0.45, 1.68, 0.42, value, 22, True, accent)


def add_notes(slide, text: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = text


def slide_notes_markdown(title: str, notes: str) -> str:
    return f"## {title}\n\n{notes}\n"


def build_ppt(data: LegacyBundle, figures: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    notes_sections: list[str] = []

    summary = data.bundle.summary
    scene_counts = pd.Series(summary["scene_counts"]).reindex(SCENE_ORDER)
    dominant_scene = scene_counts.idxmax()
    dominant_scene_count = int(scene_counts.max())
    scene_ranked = pd.Series(summary["scene_counts"]).sort_values(ascending=False)
    top2_scene_share = scene_ranked.head(2).sum() / scene_ranked.sum() if scene_ranked.sum() else 0
    scene_profiles = data.scene_profiles.sort_values("case_count", ascending=False)
    scene_product_pct = data.scene_product.div(data.scene_product.sum(axis=1).replace(0, 1), axis=0).mul(100)
    province_counts = data.cases["省份_标准"].dropna().value_counts()
    top_provinces = data.cases["省份_标准"].dropna().value_counts().head(4)
    top4_province_share = top_provinces.sum() / summary["case_count"] if summary["case_count"] else 0
    east_cases = int((data.cases["区域"] == "东部").sum())
    middle_cases = int((data.cases["区域"] == "中部").sum())
    west_cases = int((data.cases["区域"] == "西部").sum())
    east_top = data.region_product.loc["东部"].sort_values(ascending=False)
    middle_top = data.region_product.loc["中部"].sort_values(ascending=False)
    west_top = data.region_product.loc["西部"].sort_values(ascending=False)
    east_top3_share = east_top.head(3).sum() / east_top.sum() if east_top.sum() else 0
    middle_top3_share = middle_top.head(3).sum() / middle_top.sum() if middle_top.sum() else 0
    west_top3_share = west_top.head(3).sum() / west_top.sum() if west_top.sum() else 0
    subject_scene_pct = pd.crosstab(data.cases["学科"].fillna("未提及"), data.cases["应用场景（一级）"], normalize="columns").mul(100)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PPT_THEME["bg"])
    add_text(slide, 0.62, 0.48, 6.8, 0.6, "Section 2.2 教育应用现状综述", 28, True)
    add_text(slide, 0.64, 1.06, 7.2, 0.52, "当前 Word 2.2 结论同步升级为公开宣讲版：观点更锋利，证据更集中，落点更适合大型场合表达", 14, False, PPT_THEME["muted"])
    add_metric_card(slide, 0.72, 1.9, "案例数", str(summary["case_count"]), PPT_THEME["accent"])
    add_metric_card(slide, 2.9, 1.9, "工具记录", str(summary["tool_rows"]), PPT_THEME["olive"])
    add_metric_card(slide, 5.08, 1.9, "标准化产品", str(summary["product_count"]), PPT_THEME["accent"])
    add_metric_card(slide, 7.26, 1.9, "覆盖省份", str(summary["province_count"]), PPT_THEME["olive"])
    add_panel(slide, 0.72, 3.34, 6.95, 2.88)
    add_text(slide, 1.0, 3.64, 6.2, 1.9, "这份汇报要回答的不是“AI 有没有进校园”，而是三件更重要的事：\n1. 为什么七成以上案例仍集中在助学，但我们已经不能只讲助学？\n2. 为什么区域差异的本质，不是有没有接入，而是谁先形成了稳定的产品组合？\n3. 为什么下一阶段真正的竞争，不在模型热度，而在教育场景控制力？", 17)
    slide.shapes.add_picture(str(LEGACY_FIGURES["scenario"]), PptInches(8.2), PptInches(1.52), width=PptInches(4.35))
    cover_notes = (
        f"这一页先交代口径和结论框架。当前 V6 一共包含 {summary['case_count']} 个去重案例、{summary['tool_rows']} 条工具记录。"
        f"我的核心判断是，AI 教育应用已经不是简单的助学独大，而是进入了场景分化、区域分化和产品分化同步发生的新阶段。"
    )
    add_notes(slide, cover_notes)
    notes_sections.append(slide_notes_markdown("Slide 1 封面与结论框架", cover_notes))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PPT_THEME["bg"])
    add_text(slide, 0.68, 0.4, 8.4, 0.5, "一、真正的变化不是场景变多，而是主航道已经收敛出来", 24, True)
    add_text(slide, 0.7, 0.92, 8.2, 0.32, "横轴看覆盖省份数，纵轴看涉及学科数，气泡大小看场景案例规模", 12, False, PPT_THEME["muted"])
    add_panel(slide, 0.64, 1.32, 7.2, 5.55)
    slide.shapes.add_picture(str(figures["scene_bubble"]), PptInches(0.86), PptInches(1.58), width=PptInches(6.75))
    add_panel(slide, 8.14, 1.32, 4.45, 5.55)
    add_text(slide, 8.44, 1.62, 3.85, 1.75, f"助学单一场景就占到 {pct(dominant_scene_count, scene_ranked.sum())}，而助学+助教合计达到 {pct(top2_scene_share, 1)}。这说明行业已经不是“六路并进”，而是从泛试点走向主航道收敛。", 16)
    add_text(slide, 8.44, 3.62, 3.7, 2.05, f"但收敛不等于单一。助教、助评、助育虽然体量小，却在省份覆盖与学科覆盖上持续抬升，意味着更高复杂度的教学系统场景正在进入可复制阶段。", 16)
    bubble_notes = (
        f"这张图最值得讲的，不是 {dominant_scene} 最大，而是它虽然占了 {pct(dominant_scene_count, scene_ranked.sum())}，却已经不能代表全部趋势。"
        f"前两大场景合计 {pct(top2_scene_share, 1)}，说明行业正在出现主航道；与此同时，助评、助育等更复杂场景在覆盖广度上抬升，说明 AI 教育应用开始从学生支持走向教学系统深处。"
    )
    add_notes(slide, bubble_notes)
    notes_sections.append(slide_notes_markdown("Slide 2 主航道收敛", bubble_notes))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PPT_THEME["bg"])
    add_text(slide, 0.68, 0.4, 8.6, 0.5, "二、区域空间格局已经呈现“东部生态化、中部平台化、西部导入期”", 24, True)
    add_text(slide, 0.7, 0.92, 9.4, 0.32, "直接复用 Word 2.2.2 的当前论证链：省域分布地图负责看总量，区域×产品热力图负责看组合结构", 12, False, PPT_THEME["muted"])
    add_panel(slide, 0.64, 1.3, 6.0, 2.3)
    slide.shapes.add_picture(str(LEGACY_FIGURES["province"]), PptInches(0.88), PptInches(1.54), width=PptInches(5.5))
    add_panel(slide, 0.64, 3.9, 6.0, 2.95)
    slide.shapes.add_picture(str(LEGACY_FIGURES["region_product"]), PptInches(0.88), PptInches(4.12), width=PptInches(5.5))
    add_panel(slide, 6.92, 1.3, 5.66, 5.55)
    add_text(slide, 7.22, 1.66, 5.02, 0.72, "这页不再谈抽象的“结构长什么样”，而是直接给出当前 Word 已经成立的三条证据链。", 16, True, PPT_THEME["accent"])
    add_text(slide, 7.22, 2.32, 5.0, 1.18, f"1. 头部板块已经非常集中。{list_phrase(top_provinces.index.tolist(), 4)}四省市合计占全部案例的{pct(top4_province_share, 1)}，这意味着全国应用并不是均匀铺开，而是由少数头部区域牵引。", 15)
    add_text(slide, 7.22, 3.76, 5.0, 1.02, f"2. 东部的领先不是“多一点”，而是已经进入生态化组合。东部样本 {east_cases} 例，且{top_items_text(east_top, 3)}三者就构成了{pct(east_top3_share, 1)}的头部产品集中度。", 15)
    add_text(slide, 7.22, 5.0, 5.0, 1.34, f"3. 中西部的差别也很清楚。中部更像平台化稳步推进，{top_items_text(middle_top, 3)}三者集中度为{pct(middle_top3_share, 1)}；西部则更像头部产品导入期，{top_items_text(west_top, 3)}三者集中度达到{pct(west_top3_share, 1)}。", 15)
    region_notes = (
        f"这一页要讲清三层意思。第一，{list_phrase(top_provinces.index.tolist(), 4)}四省市合计占比达到{pct(top4_province_share, 1)}，头部区域板块已经成形。"
        f"第二，东部不只是案例多，而是{top_items_text(east_top, 3)}形成了更高阶的组合结构；第三，中部更像平台化推进，西部更像头部产品导入期。"
        "所以区域差异的本质不是有没有接入 AI，而是谁先形成了稳定、可复制、可扩展的产品组合。"
    )
    add_notes(slide, region_notes)
    notes_sections.append(slide_notes_markdown("Slide 3 区域空间格局证据", region_notes))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PPT_THEME["bg"])
    add_text(slide, 0.68, 0.4, 8.6, 0.5, "三、场景不是平均落地，而是深度绑定不同教育单元", 24, True)
    add_text(slide, 0.7, 0.92, 8.2, 0.32, "按场景内部归一化后，能看到它们最依赖哪些学段", 12, False, PPT_THEME["muted"])
    add_panel(slide, 0.64, 1.32, 6.5, 5.55)
    slide.shapes.add_picture(str(figures["stage_scene"]), PptInches(0.86), PptInches(1.56), width=PptInches(6.04))
    add_panel(slide, 7.42, 1.32, 5.16, 5.55)
    insight_lines = []
    for scene in SCENE_ORDER[:4]:
        row = data.scene_profiles.loc[scene]
        top_stage_pct = float(data.stage_scene_pct.loc[row['top_stage'], scene]) if row['top_stage'] in data.stage_scene_pct.index else 0.0
        top_subject_pct = float(subject_scene_pct.loc[row['top_subject'], scene]) if row['top_subject'] in subject_scene_pct.index else 0.0
        insight_lines.append(f"{scene}：{row['top_stage']}占{top_stage_pct:.0f}%，{row['top_subject']}占{top_subject_pct:.0f}%")
    add_text(slide, 7.74, 1.7, 4.54, 3.7, "\n".join(insight_lines), 16)
    add_text(slide, 7.74, 5.38, 4.3, 1.1, "这页真正要讲的是：同一模型进入不同场景后，决定它效果上限的往往不是模型参数，而是它嵌入了哪个教育组织单元。", 15, True, PPT_THEME["accent"])
    stage_notes = (
        "这一页强调的不是哪个学段多，而是每个场景在依附哪个教育单元。"
        "如果一个场景对某一学段和某一学科的依赖特别高，就说明它还处在强绑定阶段；只有这种依赖逐步下降，这个场景才可能真正泛化。"
    )
    add_notes(slide, stage_notes)
    notes_sections.append(slide_notes_markdown("Slide 4 场景与教育单元依附关系", stage_notes))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PPT_THEME["bg"])
    add_text(slide, 0.68, 0.4, 8.2, 0.5, "四、产品竞争已经从“谁更热”升级为“谁能控制关键场景”", 24, True)
    add_text(slide, 0.7, 0.92, 8.4, 0.32, "热力值不是总频次，而是场景内部占比，因此更能看出“谁在主导这个场景”", 12, False, PPT_THEME["muted"])
    add_panel(slide, 0.64, 1.32, 8.05, 5.55)
    slide.shapes.add_picture(str(figures["scene_product"]), PptInches(0.86), PptInches(1.56), width=PptInches(7.58))
    add_panel(slide, 8.94, 1.32, 3.64, 5.55)
    weakest_scene = scene_product_pct.max(axis=1).sort_values().index[0]
    weakest_pct = scene_product_pct.max(axis=1).sort_values().iloc[0]
    add_text(slide, 9.2, 1.72, 3.1, 2.0, f"这一页最重要的不是谁排第一，而是控制力已经开始分层。最强单点仍在{dominant_scene}相关场景，但也有场景的头部占比只有约{weakest_pct:.0f}%，说明仍处开放竞争状态。", 16)
    strongest_scene = scene_product_pct.max(axis=1).sort_values(ascending=False).index[0]
    strongest_product = scene_product_pct.loc[strongest_scene].sort_values(ascending=False).index[0]
    strongest_pct = scene_product_pct.loc[strongest_scene].sort_values(ascending=False).iloc[0]
    add_text(slide, 9.2, 4.18, 3.1, 1.28, f"最强单点：{strongest_scene}场景中的{strongest_product}，场景内占比约{strongest_pct:.0f}%\n最开放场景：{weakest_scene}，头部占比约{weakest_pct:.0f}%", 15, True, PPT_THEME["accent"])
    product_notes = (
        "最后这一页要讲的不是品牌热度，而是控制力分化。"
        f"目前最强单点关系出现在 {strongest_scene} 场景中的 {strongest_product}，占比约 {strongest_pct:.0f}%；但 {weakest_scene} 这样的场景头部占比只有约 {weakest_pct:.0f}%，仍是开放竞争。"
        "这意味着真正的胜负手不是谁更红，而是谁先在关键教育场景中形成不可替代的位置。"
    )
    add_notes(slide, product_notes)
    notes_sections.append(slide_notes_markdown("Slide 5 场景控制力竞争", product_notes))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PPT_THEME["bg"])
    add_text(slide, 0.68, 0.52, 6.4, 0.48, "五、汇报结论：下一阶段看的是结构升级，而不是简单铺量", 24, True)
    add_panel(slide, 0.72, 1.46, 12.0, 4.9)
    add_text(slide, 1.02, 1.88, 11.25, 3.45, "1. 从场景看：教育 AI 已从泛试点进入主航道收敛阶段，助学与助教构成最清晰的两条主线。\n\n2. 从区域看：全国格局并非均匀扩散，而是少数头部区域先形成产品组合与应用路径，东部生态化、中部平台化、西部导入期的分层已经出现。\n\n3. 从产品看：下一阶段真正决定胜负的，不是模型热度，而是谁能在关键教育场景里形成稳定控制力。", 22)
    closing_notes = (
        "最后要把判断收住。接下来不应再只看案例数增长，而要看三件事：主航道之外的新场景能否继续抬升，中西部能否从导入走向稳定的平台化与生态化，以及头部产品能否在关键教育场景中建立长期控制力。"
    )
    add_notes(slide, closing_notes)
    notes_sections.append(slide_notes_markdown("Slide 6 汇报结论", closing_notes))

    prs.save(PPTX_PATH)
    NOTES_PATH.write_text("# Section 2.2 汇报讲稿\n\n" + "\n".join(notes_sections), encoding="utf-8")


def main() -> None:
    configure_matplotlib()
    ensure_legacy_figures()
    data = build_legacy_bundle()
    regenerate_legacy_product_heatmaps(data)
    ppt_figures = build_ppt_figures(data)
    build_word_doc(data)
    build_ppt(data, ppt_figures)
    print(f"Generated docx: {DOCX_PATH}")
    print(f"Generated pptx: {PPTX_PATH}")
    print(f"Generated notes: {NOTES_PATH}")


if __name__ == "__main__":
    main()
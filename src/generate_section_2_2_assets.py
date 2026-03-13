from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
import re
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd
import seaborn as sns
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor, Cm
from matplotlib import font_manager
from matplotlib.colors import LinearSegmentedColormap
import squarify
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt, Emu

BASE_DIR = Path(__file__).resolve().parent.parent
WORKBOOK_PATH = BASE_DIR / "data" / "V6.xlsx"
DOCX_PATH = BASE_DIR / "Section_2.2_应用现状综述_最终版.docx"
PPTX_PATH = BASE_DIR / "Section_2.2_应用现状综述_汇报版.pptx"
FIG_DIR = BASE_DIR / "output" / "figures" / "section_2_2_v6"

SCENE_ORDER = ["助学", "助教", "助评", "助育", "助管", "助研"]
SCENE_COLORS = {
    "助学": "#3B82F6", "助教": "#F59E0B", "助评": "#EF4444",
    "助育": "#10B981", "助管": "#8B5CF6", "助研": "#EC4899",
}
SCENE_PALETTE = list(SCENE_COLORS.values())

ACCENT1 = "#3B82F6"
ACCENT2 = "#F59E0B"
ACCENT3 = "#10B981"
BG_DARK = "#1E293B"
BG_CARD = "#FFFFFF"
TEXT_DARK = "#1E293B"
TEXT_MID = "#64748B"
GRADIENT_CMAP = LinearSegmentedColormap.from_list("custom_heat", ["#EFF6FF", "#3B82F6", "#1E3A8A"])

FONT_FAMILY = "Microsoft YaHei"


def configure_matplotlib() -> None:
    import os
    # 1) 清除 matplotlib 字体缓存，避免缓存不一致导致中文乱码
    cache_dir = matplotlib.get_cachedir()
    if cache_dir and os.path.isdir(cache_dir):
        for fname in os.listdir(cache_dir):
            if "font" in fname.lower():
                try:
                    os.remove(os.path.join(cache_dir, fname))
                except OSError:
                    pass
    # 2) 重建字体管理器
    font_manager._load_fontmanager(try_read_cache=False)

    candidates = ["Microsoft YaHei", "SimHei", "SimSun", "KaiTi", "Noto Sans SC", "DejaVu Sans"]
    available = {f.name for f in font_manager.fontManager.ttflist}
    global FONT_FAMILY
    FONT_FAMILY = next((n for n in candidates if n in available), "DejaVu Sans")
    print(f"[字体] 使用: {FONT_FAMILY}")

    # 3) 先调用 seaborn set_theme（它会重置 rcParams）
    sns.set_theme(style="whitegrid")

    # 4) 在 seaborn 之后设置字体，确保不被覆盖
    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.sans-serif": [FONT_FAMILY, "SimHei", "DejaVu Sans"],
        "axes.unicode_minus": False,
        "figure.dpi": 150,
        "savefig.dpi": 300,
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "axes.edgecolor": "#E2E8F0",
        "axes.grid": True,
        "grid.color": "#F1F5F9",
        "grid.linewidth": 0.8,
    })


def normalize_text(value) -> str | None:
    if pd.isna(value):
        return None
    text = str(value).strip()
    if text in {"", "nan", "None", "未提及", "未知"}:
        return None
    return text


def normalize_product_name(value) -> str | None:
    text = normalize_text(value)
    if text is None:
        return None
    text = re.sub(r"\s+", " ", str(text)).strip()
    alias = {"即梦": "即梦AI", "即梦 AI": "即梦AI", "剪映": "剪映AI",
             "剪映 AI": "剪映AI", "DeepSeek": "DeepSeek 大模型", "语文朗读宝 AI": "语文朗读宝AI"}
    return alias.get(text, text)


def normalize_stage(value) -> str:
    if pd.isna(value):
        return "未提及"
    text = str(value)
    for kw, label in [("小学/初中/高中", "全学段"), ("小学至初中", "小学/初中"),
                       ("初中/高中", "初中/高中"), ("幼儿园", "学前"), ("学前", "学前"),
                       ("小学", "小学"), ("初中", "初中"), ("高中", "高中"),
                       ("中学", "中学"), ("中职", "中职"), ("职高", "中职")]:
        if kw in text:
            return label
    return text.strip() or "未提及"


def pct(part: float, whole: float) -> str:
    return f"{part / whole * 100:.1f}%" if whole else "0.0%"


@dataclass
class AnalysisBundle:
    df: pd.DataFrame
    cases: pd.DataFrame
    stage_scene: pd.DataFrame
    subject_scene: pd.DataFrame
    product_type_scene: pd.DataFrame
    product_attr_scene: pd.DataFrame
    llm_scene: pd.DataFrame
    summary: dict = field(default_factory=dict)


def load_bundle() -> AnalysisBundle:
    df = pd.read_excel(WORKBOOK_PATH, sheet_name="（新）processed_results")
    df = df.drop(columns=[c for c in df.columns if str(c).startswith("Unnamed")], errors="ignore")
    for col in df.columns:
        if df[col].dtype == object:
            df[col] = df[col].map(normalize_text)

    df["案例编号"] = pd.to_numeric(df["案例编号"], errors="coerce")
    df["省份_标准"] = df["省份_更新"].fillna(df["省份"])
    df["产品名_标准"] = df["产品名（校准）"].fillna(df["工具标准名"]).map(normalize_product_name)

    # 核心去重：按文件名去重，一个案例只出现一次
    if "文件名" in df.columns:
        df = df.drop_duplicates(subset=["文件名"], keep="first")
        print(f"[去重] 按文件名去重后剩余 {len(df)} 条记录")

    cases = df.dropna(subset=["案例编号"]).sort_values("案例编号").drop_duplicates("案例编号", keep="first").copy()
    cases["学段_标准"] = cases["学段"].apply(normalize_stage)
    cases["产品属性_标准"] = cases["产品属性"].fillna("其他")
    cases["是否大模型_标准"] = cases["产品名_是否大模型"].fillna("否")
    cases["产品分类_标准"] = cases["产品分类"].fillna("未提及")

    # 交叉表
    stage_order = ["学前", "小学", "初中", "高中", "中学", "小学/初中", "初中/高中", "全学段", "中职", "未提及"]
    stage_scene = pd.crosstab(cases["学段_标准"], cases["应用场景（一级）"])
    stage_scene = stage_scene.reindex(index=[s for s in stage_order if s in stage_scene.index], fill_value=0)
    stage_scene = stage_scene.reindex(columns=SCENE_ORDER, fill_value=0)

    subj_counts = cases["学科"].fillna("未提及").value_counts()
    subj_top = [s for s in subj_counts.index if s != "未提及"][:10]
    mask = cases["学科"].isin(subj_top)
    subject_scene = pd.crosstab(cases.loc[mask, "学科"], cases.loc[mask, "应用场景（一级）"])
    subject_scene = subject_scene.reindex(index=subj_top, fill_value=0).reindex(columns=SCENE_ORDER, fill_value=0)

    pt_order = cases["产品分类_标准"].value_counts().index.tolist()
    product_type_scene = pd.crosstab(cases["产品分类_标准"], cases["应用场景（一级）"])
    product_type_scene = product_type_scene.reindex(index=pt_order, fill_value=0).reindex(columns=SCENE_ORDER, fill_value=0)

    pa_valid = cases[cases["产品属性_标准"].isin(["AI智能体", "大语言模型", "其他"])]
    product_attr_scene = pd.crosstab(pa_valid["产品属性_标准"], pa_valid["应用场景（一级）"])
    product_attr_scene = product_attr_scene.reindex(index=["AI智能体", "大语言模型", "其他"], fill_value=0)
    product_attr_scene = product_attr_scene.reindex(columns=SCENE_ORDER, fill_value=0)

    llm_valid = cases[cases["是否大模型_标准"].isin(["是", "否"])]
    llm_scene = pd.crosstab(llm_valid["是否大模型_标准"], llm_valid["应用场景（一级）"])
    llm_scene = llm_scene.reindex(index=["是", "否"], fill_value=0).reindex(columns=SCENE_ORDER, fill_value=0)

    l1_counts = cases["应用场景（一级）"].value_counts().reindex(SCENE_ORDER, fill_value=0)
    l2_pairs = (cases.dropna(subset=["应用场景（一级）", "应用场景（二级）"])
                .groupby(["应用场景（一级）", "应用场景（二级）"]).size()
                .reset_index(name="count").sort_values(["应用场景（一级）", "count"], ascending=[True, False]))
    top_l2 = {}
    for sc in SCENE_ORDER:
        sub = l2_pairs[l2_pairs["应用场景（一级）"] == sc].head(2)
        top_l2[sc] = [f"{r['应用场景（二级）']}（{int(r['count'])}例）" for _, r in sub.iterrows()]

    total = int(cases["案例编号"].nunique())
    summary = {
        "total_cases": total,
        "tool_rows": int(len(df)),
        "product_count": int(df["产品名_标准"].dropna().nunique()),
        "province_count": int(cases["省份_标准"].dropna().nunique()),
        "scene_counts": l1_counts.to_dict(),
        "top_l2_by_l1": top_l2,
        "stage_counts": cases["学段_标准"].value_counts().to_dict(),
        "subject_counts": cases["学科"].fillna("未提及").value_counts().to_dict(),
        "province_counts": cases["省份_标准"].fillna("未提及").value_counts().to_dict(),
        "top_products": df["产品名_标准"].dropna().value_counts().head(15).to_dict(),
    }

    return AnalysisBundle(df=df, cases=cases, stage_scene=stage_scene,
                          subject_scene=subject_scene, product_type_scene=product_type_scene,
                          product_attr_scene=product_attr_scene, llm_scene=llm_scene, summary=summary)


# ---------------------------------------------------------------------------
# Chart generation helpers
# ---------------------------------------------------------------------------

def _finish(fig, path):
    fig.tight_layout()
    fig.savefig(path, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  [图表] {path.name}")


def fig_scene_bar(bundle: AnalysisBundle, path: Path):
    """① 六大场景柱状图"""
    s = pd.Series(bundle.summary["scene_counts"]).reindex(SCENE_ORDER, fill_value=0)
    total = s.sum()
    fig, ax = plt.subplots(figsize=(10, 5.5))
    bars = ax.bar(s.index, s.values, color=SCENE_PALETTE, width=0.6, edgecolor="white", linewidth=1.2)
    for bar, v in zip(bars, s.values):
        ax.text(bar.get_x() + bar.get_width()/2, v + total*0.008,
                f"{int(v)}\n({v/total*100:.1f}%)", ha="center", va="bottom", fontsize=10, fontweight="bold")
    ax.set_title("六大一级应用场景案例分布", fontsize=18, fontweight="bold", pad=20, color=TEXT_DARK)
    ax.set_ylabel("案例数", fontsize=12)
    ax.set_ylim(0, s.max() * 1.18)
    sns.despine()
    _finish(fig, path)


def fig_scene_treemap(bundle: AnalysisBundle, path: Path):
    """② 场景 Treemap（仅一级）"""
    s = pd.Series(bundle.summary["scene_counts"]).reindex(SCENE_ORDER, fill_value=0)
    total = s.sum()
    labels = [f"{name}\n{int(v)} 例\n({v/total*100:.1f}%)" for name, v in s.items()]
    colors = [SCENE_COLORS[name] for name in s.index]
    fig, ax = plt.subplots(figsize=(11, 6.5))
    squarify.plot(sizes=s.values, label=labels, color=colors, alpha=0.88,
                  text_kwargs={"fontsize": 13, "fontweight": "bold", "color": "white"}, ax=ax)
    ax.set_title("六大应用场景占比（一级分类）", fontsize=18, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.axis("off")
    _finish(fig, path)


def fig_stage_bar(bundle: AnalysisBundle, path: Path):
    """③ 学段分布水平柱状图"""
    s = pd.Series(bundle.summary["stage_counts"])
    s = s.drop(labels=["未提及"], errors="ignore").sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(9.5, 5.5))
    colors = plt.cm.Blues(np.linspace(0.35, 0.85, len(s)))
    ax.barh(s.index, s.values, color=colors, height=0.6, edgecolor="white")
    for i, v in enumerate(s.values):
        ax.text(v + s.max()*0.015, i, f"{int(v)}", va="center", fontsize=11, fontweight="bold")
    ax.set_title("案例学段分布", fontsize=18, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_xlabel("案例数", fontsize=12)
    sns.despine()
    _finish(fig, path)


def fig_subject_lollipop(bundle: AnalysisBundle, path: Path):
    """④ 学科 Top10 棒棒糖图"""
    s = pd.Series(bundle.summary["subject_counts"]).drop(labels=["未提及"], errors="ignore").head(10)
    s = s.sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(9.5, 6))
    y = range(len(s))
    colors = plt.cm.Oranges(np.linspace(0.4, 0.85, len(s)))
    ax.hlines(y=list(y), xmin=0, xmax=s.values, color=colors, linewidth=2.5)
    ax.scatter(s.values, list(y), color=colors, s=100, zorder=5, edgecolors="white", linewidth=1.5)
    ax.set_yticks(list(y))
    ax.set_yticklabels(s.index, fontsize=11)
    for i, v in enumerate(s.values):
        ax.text(v + s.max()*0.02, i, str(int(v)), va="center", fontsize=11, fontweight="bold")
    ax.set_title("案例学科 Top 10", fontsize=18, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_xlabel("案例数", fontsize=12)
    ax.set_xlim(0, s.max()*1.15)
    sns.despine()
    _finish(fig, path)


def fig_stage_scene_heatmap(bundle: AnalysisBundle, path: Path):
    """⑤ 学段×场景 热力图"""
    df = bundle.stage_scene.copy()
    df = df.loc[df.sum(axis=1) > 0]
    fig, ax = plt.subplots(figsize=(10, max(5.5, 0.6*len(df)+2.5)))
    sns.heatmap(df, annot=True, fmt="g", cmap=GRADIENT_CMAP, linewidths=0.8,
                linecolor="white", cbar_kws={"shrink": 0.8, "label": "案例数"}, ax=ax)
    ax.set_title("学段 × 应用场景 交叉分布", fontsize=16, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_xlabel(""); ax.set_ylabel("")
    ax.tick_params(axis="both", labelsize=11)
    _finish(fig, path)


def fig_subject_scene_heatmap(bundle: AnalysisBundle, path: Path):
    """⑥ 学科×场景 热力图"""
    df = bundle.subject_scene.copy()
    fig, ax = plt.subplots(figsize=(10, max(5.5, 0.6*len(df)+2.5)))
    sns.heatmap(df, annot=True, fmt="g", cmap=GRADIENT_CMAP, linewidths=0.8,
                linecolor="white", cbar_kws={"shrink": 0.8, "label": "案例数"}, ax=ax)
    ax.set_title("Top10 学科 × 应用场景 交叉分布", fontsize=16, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_xlabel(""); ax.set_ylabel("")
    ax.tick_params(axis="both", labelsize=11)
    _finish(fig, path)


def fig_product_type_scene(bundle: AnalysisBundle, path: Path):
    """⑦ 产品分类×场景 堆叠横向条形图"""
    df = bundle.product_type_scene.loc[bundle.product_type_scene.sum(axis=1) > 0].copy()
    fig, ax = plt.subplots(figsize=(11, max(5, 0.7*len(df)+2.5)))
    left = pd.Series(0.0, index=df.index)
    for sc in SCENE_ORDER:
        if sc in df.columns:
            vals = df[sc]
            ax.barh(df.index, vals, left=left, color=SCENE_COLORS[sc], label=sc,
                    height=0.55, edgecolor="white", linewidth=0.8)
            left += vals
    ax.set_title("产品分类 × 应用场景 交叉分布", fontsize=16, fontweight="bold", pad=30, color=TEXT_DARK)
    ax.set_xlabel("案例数", fontsize=12)
    ax.legend(ncol=6, bbox_to_anchor=(0.5, 1.02), loc="lower center", frameon=False, fontsize=10)
    ax.invert_yaxis()
    sns.despine()
    _finish(fig, path)


def fig_product_attr_scene(bundle: AnalysisBundle, path: Path):
    """⑧ 产品属性×场景 分组柱状图"""
    df = bundle.product_attr_scene.copy()
    n_groups = len(SCENE_ORDER)
    n_bars = len(df)
    x = np.arange(n_groups)
    w = 0.22
    attr_colors = {"AI智能体": "#3B82F6", "大语言模型": "#F59E0B", "其他": "#94A3B8"}
    fig, ax = plt.subplots(figsize=(11, 5.5))
    for i, (attr, row) in enumerate(df.iterrows()):
        offset = (i - n_bars/2 + 0.5) * w
        bars = ax.bar(x + offset, [row.get(s, 0) for s in SCENE_ORDER], w,
                      label=attr, color=attr_colors.get(attr, "#94A3B8"),
                      edgecolor="white", linewidth=0.8)
        for bar in bars:
            h = bar.get_height()
            if h > 0:
                ax.text(bar.get_x()+bar.get_width()/2, h+2, str(int(h)),
                        ha="center", va="bottom", fontsize=8, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(SCENE_ORDER, fontsize=12)
    ax.set_title("产品属性 × 应用场景 分布", fontsize=16, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_ylabel("案例数", fontsize=12)
    ax.legend(frameon=False, fontsize=11)
    sns.despine()
    _finish(fig, path)


def fig_province_bar(bundle: AnalysisBundle, path: Path):
    """⑨ 省份 Top10"""
    s = pd.Series(bundle.summary["province_counts"]).drop(labels=["未提及"], errors="ignore").head(10)
    s = s.sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(9.5, 5.5))
    colors = plt.cm.Greens(np.linspace(0.35, 0.85, len(s)))
    ax.barh(s.index, s.values, color=colors, height=0.6, edgecolor="white")
    for i, v in enumerate(s.values):
        ax.text(v + s.max()*0.015, i, str(int(v)), va="center", fontsize=11, fontweight="bold")
    ax.set_title("案例量最高的 10 个省份", fontsize=18, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_xlabel("案例数", fontsize=12)
    sns.despine()
    _finish(fig, path)


def fig_llm_scene(bundle: AnalysisBundle, path: Path):
    """⑩ 是否大模型×场景 对比柱状图"""
    df = bundle.llm_scene.copy()
    x = np.arange(len(SCENE_ORDER))
    w = 0.32
    fig, ax = plt.subplots(figsize=(11, 5.5))
    clr = {"是": "#EF4444", "否": "#3B82F6"}
    lbl = {"是": "大模型产品", "否": "非大模型产品"}
    for i, (key, row) in enumerate(df.iterrows()):
        offset = (i - 0.5) * w
        bars = ax.bar(x + offset, [row.get(s, 0) for s in SCENE_ORDER], w,
                      label=lbl.get(key, key), color=clr.get(key, "#94A3B8"),
                      edgecolor="white", linewidth=0.8)
        for bar in bars:
            h = bar.get_height()
            if h > 0:
                ax.text(bar.get_x()+bar.get_width()/2, h+2, str(int(h)),
                        ha="center", va="bottom", fontsize=9, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(SCENE_ORDER, fontsize=12)
    ax.set_title("大模型 vs 非大模型 × 应用场景", fontsize=16, fontweight="bold", pad=15, color=TEXT_DARK)
    ax.set_ylabel("案例数", fontsize=12)
    ax.legend(frameon=False, fontsize=11)
    sns.despine()
    _finish(fig, path)


def generate_figures(bundle: AnalysisBundle) -> dict[str, Path]:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    paths = {}
    chart_funcs = [
        ("scene_bar", fig_scene_bar),
        ("scene_treemap", fig_scene_treemap),
        ("stage_bar", fig_stage_bar),
        ("subject_lollipop", fig_subject_lollipop),
        ("stage_scene_heatmap", fig_stage_scene_heatmap),
        ("subject_scene_heatmap", fig_subject_scene_heatmap),
        ("product_type_scene", fig_product_type_scene),
        ("product_attr_scene", fig_product_attr_scene),
        ("province_bar", fig_province_bar),
        ("llm_scene", fig_llm_scene),
    ]
    for name, func in chart_funcs:
        p = FIG_DIR / f"fig_s22_{name}.png"
        func(bundle, p)
        paths[name] = p
    return paths


# ---------------------------------------------------------------------------
# DOCX builder
# ---------------------------------------------------------------------------

def _set_docx_style(doc: Document):
    style = doc.styles["Normal"]
    style.font.name = "宋体"
    style._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)


def _add_para(doc, text, bold_prefix=None):
    p = doc.add_paragraph()
    p.paragraph_format.line_spacing = Pt(22)
    p.paragraph_format.space_after = Pt(6)
    if bold_prefix and text.startswith(bold_prefix):
        r = p.add_run(bold_prefix); r.bold = True
        p.add_run(text[len(bold_prefix):])
    else:
        p.add_run(text)


def _add_caption(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    r.font.size = Pt(10); r.bold = True; r.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def _insert_pic(doc, path, width=5.8):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Inches(width))


def build_docx(bundle: AnalysisBundle, figs: dict[str, Path]):
    doc = Document()
    _set_docx_style(doc)
    s = bundle.summary
    total = s["total_cases"]
    sc = pd.Series(s["scene_counts"]).reindex(SCENE_ORDER, fill_value=0)
    top_scene, top_scene_n = sc.idxmax(), int(sc.max())
    st = s["stage_counts"]
    top_stage = next(iter(st)); top_stage_n = st[top_stage]
    sj = {k: v for k, v in s["subject_counts"].items() if k != "未提及"}
    top_subj = next(iter(sj)); top_subj_n = sj[top_subj]
    prov = {k: v for k, v in s["province_counts"].items() if k != "未提及"}
    top_prov = next(iter(prov)); top_prov_n = prov[top_prov]

    # 2.2 总览
    doc.add_heading("2.2 教育应用类产品应用现状综述", level=1)
    _add_para(doc, f'基于 V6.xlsx 全量更新（主事实表经文件名去重后共 {total} 个独立案例、{s["product_count"]} 个标准化产品、覆盖 {s["province_count"]} 个省级单元），本节从学段学科、应用场景、产品形态和技术路径四个维度，系统梳理 AI 教育应用的实践格局。')

    # 2.2.1
    doc.add_heading("2.2.1 数据基础与工作簿理解", level=2)
    sheet_desc = (
        'V6 工作簿已形成"主事实表 + 场景分类字典 + 产品分类参考 + 头部产品说明 + 省份映射"'
        '的完整分析结构。主事实表"（新）processed_results"原始记录 '
        f'{s["tool_rows"]} 条工具级记录，经按文件名去重后回收为 {total} 个独立案例。'
        '本节所有分析均基于去重后的案例级口径。'
    )
    _add_para(doc, sheet_desc)

    # 2.2.2
    doc.add_heading("2.2.2 学段学科角度：应用落地的主阵地", level=2)
    _add_para(doc, f'从案例级口径看，{top_stage}是 AI 教育应用最集中的落地学段，共 {top_stage_n} 例（{pct(top_stage_n, total)}）。初中、高中紧随其后，说明实践主要围绕义务教育与升学关键阶段展开。跨学段案例虽已出现但体量仍小于单一学段，反映出产品适配仍以具体教学组织单元为主。')
    _add_para(doc, f'{top_subj}以 {top_subj_n} 例位居首位，语文、英语、科学等学科随后跟进，表明 AI 应用最容易切入的是既有高频教学任务、又便于形成标准化反馈的数据密集型学科。美术、体育等素养型领域在特定场景中已形成较强特色，说明正在从"主科增效"向"多学科延展"过渡。')
    _insert_pic(doc, figs["stage_bar"])
    _add_caption(doc, "图2-1 案例学段分布")
    _insert_pic(doc, figs["subject_lollipop"])
    _add_caption(doc, "图2-2 案例学科 Top 10")

    # 2.2.3
    doc.add_heading("2.2.3 场景角度：一级场景结构与占比", level=2)
    _add_para(doc, f'{top_scene}是最核心的应用方向（{top_scene_n} 例，{pct(top_scene_n, total)}），其后依次是助教、助评、助育、助管、助研，整体呈现"以学生学习支持为主轴、以教师赋能和评价改革为次主线、以治理与研究为长尾"的结构特征。')
    tbl = doc.add_table(rows=1, cols=4, style="Table Grid")
    for i, h in enumerate(["一级场景", "案例数", "占比", "核心二级场景"]):
        tbl.rows[0].cells[i].text = h
    for scene in SCENE_ORDER:
        cnt = int(sc[scene])
        row = tbl.add_row().cells
        row[0].text = scene; row[1].text = str(cnt)
        row[2].text = pct(cnt, total)
        row[3].text = "；".join(s["top_l2_by_l1"].get(scene, [])) or "-"
    _insert_pic(doc, figs["scene_bar"])
    _add_caption(doc, "图2-3 六大一级应用场景分布")
    _insert_pic(doc, figs["scene_treemap"])
    _add_caption(doc, "图2-4 六大应用场景占比（Treemap）")

    # 2.2.4
    doc.add_heading("2.2.4 多因素交叉分析", level=2)
    _add_para(doc, "为揭示不同维度之间的互动模式，本节对学段×场景、学科×场景、产品分类×场景、产品属性×场景、大模型渗透×场景进行交叉透视分析。")

    doc.add_heading("（1）学段 × 场景", level=3)
    dom_stage = {sc_name: bundle.stage_scene[sc_name].idxmax() for sc_name in SCENE_ORDER if bundle.stage_scene[sc_name].sum() > 0}
    stage_mapping = '、'.join(f'{k}\u2192{v}' for k, v in dom_stage.items())
    _add_para(doc, f'热力图显示，助学在小学学段最为集中，助教在小学和初中均有大量案例，助评和助育则在小学和幼儿园阶段呈现较高活跃度。各场景的主力学段分别为：{stage_mapping}。')
    _insert_pic(doc, figs["stage_scene_heatmap"])
    _add_caption(doc, "图2-5 学段 × 应用场景 交叉热力图")

    doc.add_heading("（2）学科 × 场景", level=3)
    dom_subj = {sc_name: bundle.subject_scene[sc_name].idxmax() for sc_name in SCENE_ORDER if bundle.subject_scene[sc_name].sum() > 0}
    subj_mapping = '、'.join(f'{k}\u2192{v}' for k, v in dom_subj.items())
    _add_para(doc, f'Top10 学科与场景的交叉分布进一步揭示了学科偏好差异：{subj_mapping}。数学、语文、英语在助学场景中占据主体，美术和体育则在助教场景呈现特色。')
    _insert_pic(doc, figs["subject_scene_heatmap"])
    _add_caption(doc, "图2-6 Top10 学科 × 应用场景 交叉热力图")

    doc.add_heading("（3）产品分类 × 场景", level=3)
    _add_para(doc, '产品分类维度显示，平台型综合系统在所有场景中均占据主体地位，Web 教学平台（SaaS）是第二大产品形态。教学辅助工具在助学场景中有较高渗透率，AI 能力/模型/API 层产品虽然总量有限但在助学和助教中已有明显应用。')
    _insert_pic(doc, figs["product_type_scene"])
    _add_caption(doc, "图2-7 产品分类 × 应用场景 交叉分布")

    doc.add_heading("（4）产品属性 × 场景", level=3)
    _add_para(doc, "AI 智能体是最主要的产品属性类型，在所有六大场景中均居首位。大语言模型型产品在助学场景的渗透率最高，反映出生成式 AI 正在成为学习辅助的核心技术路径。")
    _insert_pic(doc, figs["product_attr_scene"])
    _add_caption(doc, "图2-8 产品属性 × 应用场景 分布")

    doc.add_heading("（5）大模型渗透 × 场景", level=3)
    llm_yes = bundle.llm_scene.loc["是"].sum() if "是" in bundle.llm_scene.index else 0
    llm_total = bundle.llm_scene.sum().sum()
    _add_para(doc, f'大模型产品在全部场景中的渗透率为 {pct(llm_yes, llm_total)}。从场景分布看，助学是大模型渗透最高的场景，助研和助管中大模型占比相对较低，反映出大模型应用目前仍以直接面向师生的教学场景为主。')
    _insert_pic(doc, figs["llm_scene"])
    _add_caption(doc, "图2-9 大模型 vs 非大模型 × 应用场景")

    # 2.2.5
    doc.add_heading("2.2.5 省份分布", level=2)
    _add_para(doc, f"按省份统计，{top_prov}以 {top_prov_n} 例位列首位，头部省份以平台建设、课程整合和区域推进三种方式同步放大案例产出。")
    _insert_pic(doc, figs["province_bar"])
    _add_caption(doc, "图2-10 案例量最高的 10 个省份")

    # 2.2.6
    doc.add_heading("2.2.6 综合判断：从单点应用转向结构分化", level=2)
    _add_para(doc, "综合学段、学科、场景、产品形态和技术路径五个维度，当前 AI 教育应用可概括为三个趋势：")
    _add_para(doc, "第一，应用主阵地仍在小学与主干学科，但多学科、多学段扩散趋势已经显现。")
    _add_para(doc, '第二，六大场景已从"助学独大"向"多场景并行"演化，助教、助评、助育正在成为新的增长极。')
    _add_para(doc, "第三，产品形态以平台型和 SaaS 为主体，AI 智能体和大模型正在加速渗透，特别是在助学场景中大模型已形成显著份额，预示着生成式 AI 将深刻重塑未来教育应用的技术路径。")

    doc.save(DOCX_PATH)
    print(f"[DOCX] 已生成 {DOCX_PATH}")


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

def _add_tb(slide, left, top, width, height, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = PptPt(size); p.font.bold = bold; p.font.name = "Microsoft YaHei"
    p.alignment = align
    if color:
        p.font.color.rgb = PptRGB(*color)


def _add_pic(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width: kwargs["width"] = PptInches(width)
    if height: kwargs["height"] = PptInches(height)
    slide.shapes.add_picture(str(path), PptInches(left), PptInches(top), **kwargs)


def build_pptx(bundle: AnalysisBundle, figs: dict[str, Path]):
    prs = Presentation()
    prs.slide_width = PptInches(13.333); prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    s = bundle.summary
    total = s["total_cases"]
    sc = pd.Series(s["scene_counts"]).reindex(SCENE_ORDER, fill_value=0)
    st = s["stage_counts"]
    top_st = next(iter(st))
    sj = {k: v for k, v in s["subject_counts"].items() if k != "未提及"}
    top_sj = next(iter(sj))
    prov = {k: v for k, v in s["province_counts"].items() if k != "未提及"}
    top_prov = next(iter(prov)); top_prov_n = prov[top_prov]
    llm_yes = bundle.llm_scene.loc["是"].sum() if "是" in bundle.llm_scene.index else 0
    llm_all = bundle.llm_scene.sum().sum()

    # Slide 1: 封面
    sl = prs.slides.add_slide(blank)
    _add_tb(sl, 1.0, 1.5, 11.3, 1.2, "2.2 教育应用类产品应用现状综述", 36, True, (0x1E, 0x29, 0x3B), PP_ALIGN.CENTER)
    _add_tb(sl, 1.0, 3.0, 11.3, 0.8,
            f"数据口径：V6.xlsx · 按文件名去重后 {total} 个独立案例 · {s['product_count']} 个标准化产品 · 覆盖 {s['province_count']} 个省份",
            14, False, (0x64, 0x74, 0x8B), PP_ALIGN.CENTER)

    # Slide 2: 场景概览
    sl = prs.slides.add_slide(blank)
    _add_tb(sl, 0.4, 0.15, 12.5, 0.5, "六大应用场景概览", 26, True, (0x1E, 0x29, 0x3B))
    _add_pic(sl, figs["scene_bar"], 0.3, 0.8, width=6.8)
    _add_pic(sl, figs["scene_treemap"], 7.2, 0.8, width=5.8)
    insight = (f"• {sc.idxmax()}场景占比 {pct(int(sc.max()), total)}，是最核心应用方向\n"
               f"• 助教（{int(sc['助教'])}例）、助评（{int(sc['助评'])}例）构成第二梯队\n"
               f"• 助育、助管、助研等长尾场景已形成明确功能切口")
    _add_tb(sl, 0.5, 6.5, 12.0, 0.9, insight, 11, False, (0x47, 0x55, 0x69))

    # Slide 3: 学段与学科
    sl = prs.slides.add_slide(blank)
    _add_tb(sl, 0.4, 0.15, 12.5, 0.5, "学段与学科分布", 26, True, (0x1E, 0x29, 0x3B))
    _add_pic(sl, figs["stage_bar"], 0.3, 0.75, width=6.3)
    _add_pic(sl, figs["subject_lollipop"], 6.8, 0.75, width=6.2)
    st = s["stage_counts"]; top_st = next(iter(st))
    sj = {k: v for k, v in s["subject_counts"].items() if k != "未提及"}; top_sj = next(iter(sj))
    _add_tb(sl, 0.5, 6.5, 12.0, 0.9,
            f"• {top_st}是最主要落地学段（{st[top_st]}例） • {top_sj}是最高频学科（{sj[top_sj]}例） • 多学科扩散趋势已显现",
            11, False, (0x47, 0x55, 0x69))

    # Slide 4: 多因素交叉 1
    sl = prs.slides.add_slide(blank)
    _add_tb(sl, 0.4, 0.15, 12.5, 0.5, "多因素交叉分析：学段×场景 & 学科×场景", 24, True, (0x1E, 0x29, 0x3B))
    _add_pic(sl, figs["stage_scene_heatmap"], 0.3, 0.75, width=6.3)
    _add_pic(sl, figs["subject_scene_heatmap"], 6.8, 0.75, width=6.2)
    dom_st = {k: bundle.stage_scene[k].idxmax() for k in SCENE_ORDER if bundle.stage_scene[k].sum() > 0}
    lines = [f"• {k}场景集中在{v}学段" for k, v in list(dom_st.items())[:3]]
    _add_tb(sl, 0.5, 6.5, 12.0, 0.9, " ".join(lines), 11, False, (0x47, 0x55, 0x69))

    # Slide 5: 多因素交叉 2
    sl = prs.slides.add_slide(blank)
    _add_tb(sl, 0.4, 0.15, 12.5, 0.5, "产品分类 × 场景 & 产品属性 × 场景 & 大模型对比", 22, True, (0x1E, 0x29, 0x3B))
    _add_pic(sl, figs["product_type_scene"], 0.2, 0.75, width=4.3)
    _add_pic(sl, figs["product_attr_scene"], 4.6, 0.75, width=4.3)
    _add_pic(sl, figs["llm_scene"], 9.0, 0.75, width=4.1)
    llm_yes = bundle.llm_scene.loc["是"].sum() if "是" in bundle.llm_scene.index else 0
    llm_all = bundle.llm_scene.sum().sum()
    _add_tb(sl, 0.5, 6.5, 12.0, 0.9,
            f"• 平台型系统主导所有场景 • AI智能体占据最大份额 • 大模型渗透率 {pct(llm_yes, llm_all)}，助学场景最高",
            11, False, (0x47, 0x55, 0x69))

    # Slide 6: 综合判断
    sl = prs.slides.add_slide(blank)
    _add_tb(sl, 0.4, 0.15, 12.5, 0.5, "综合判断", 28, True, (0x1E, 0x29, 0x3B))
    _add_pic(sl, figs["province_bar"], 0.3, 0.75, width=6.0)
    conclusions = (
        "核心洞见\n\n"
        f"1. 应用主阵地：{top_st}学段 + {top_sj}学科，但多学科扩散趋势已现\n\n"
        f'2. 场景分化：从"助学独大"（{pct(int(sc.max()), total)}）向助教、助评、助育多线并行演化\n\n'
        f"3. 产品形态：平台型 + SaaS 是主体，AI智能体是最主要产品属性\n\n"
        f"4. 技术路径：大模型渗透率 {pct(llm_yes, llm_all)}，在助学场景最为突出\n\n"
        f"5. 地域格局：{top_prov}以 {top_prov_n} 例领跑，头部省份持续扩大应用规模"
    )
    _add_tb(sl, 6.5, 0.75, 6.5, 5.8, conclusions, 14, False, (0x1E, 0x29, 0x3B))

    prs.save(PPTX_PATH)
    print(f"[PPTX] 已生成 {PPTX_PATH}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    configure_matplotlib()
    print("=" * 60)
    print("Section 2.2 资产生成：开始")
    print("=" * 60)
    bundle = load_bundle()
    figs = generate_figures(bundle)
    build_docx(bundle, figs)
    build_pptx(bundle, figs)
    print("=" * 60)
    print("全部完成！")
    print(f"  DOCX: {DOCX_PATH}")
    print(f"  PPTX: {PPTX_PATH}")
    print(f"  图表: {FIG_DIR}")
    print("=" * 60)


if __name__ == "__main__":
    main()